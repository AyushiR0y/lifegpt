"""
document_parser.py
Universal document parser supporting:
  PDF, DOCX, TXT, CSV, XLSX/XLS, PPTX, MSG (Outlook e-mail)

Key dataclasses
---------------
DocumentContent  – normalised output returned by every parser
EmailChain       – structured representation of an MSG e-mail thread
"""

import io
import os
import re
import base64
import tempfile
import zipfile
from html import unescape
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Any


# ---------------------------------------------------------------------------
# Data-classes
# ---------------------------------------------------------------------------

@dataclass
class ExtractedImage:
    image_base64: str
    image_format: str
    page_number: int
    image_index: int
    description: str = ""
    width: int = 0
    height: int = 0


@dataclass
class DocumentContent:
    text: str
    total_pages: int
    page_contents: Dict[int, str]
    metadata: Dict[str, Any]
    file_type: str
    images: List[ExtractedImage] = field(default_factory=list)


@dataclass
class EmailAttachment:
    filename: str
    content_type: str
    size_bytes: int
    content: Optional[str] = None
    is_readable: bool = False


@dataclass
class ExtractedLink:
    url: str
    context: str
    link_type: str


@dataclass
class ExtractedTimeline:
    original_text: str
    context: str
    event_type: str


@dataclass
class EmailMessage:
    sender: str
    recipients: List[str]
    cc: List[str]
    subject: str
    date: str
    body: str
    attachments: List[EmailAttachment] = field(default_factory=list)
    links: List[ExtractedLink] = field(default_factory=list)
    timelines: List[ExtractedTimeline] = field(default_factory=list)
    action_requested: str = ""
    email_number: int = 0


@dataclass
class EmailChain:
    messages: List[EmailMessage]
    total_attachments: int
    attachment_details: List[EmailAttachment]
    subject: str
    participants: List[str]
    all_links: List[ExtractedLink] = field(default_factory=list)
    all_timelines: List[ExtractedTimeline] = field(default_factory=list)
    conversation_flow: str = ""


# ---------------------------------------------------------------------------
# Parser
# ---------------------------------------------------------------------------

class DocumentParser:
    SUPPORTED_FORMATS = {
        "txt": "Text File",
        "csv": "CSV File",
        "xlsx": "Excel Workbook",
        "xls": "Excel Workbook (Legacy)",
        "pdf": "PDF Document",
        "docx": "Word Document",
        "pptx": "PowerPoint Presentation",
        "msg": "Outlook Email",
    }

    MAX_IMAGES = 20
    MAX_IMAGE_DIMENSION = 1024

    _URL_RE = re.compile(r"https?://[^\s<>\"{}|\\^`\[\]]+|www\.[^\s<>\"{}|\\^`\[\]]+")

    _DATE_PATTERNS = [
        r"\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\s+\d{1,2}:\d{2}\s*(?:AM|PM|am|pm)?",
        r"\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{2,4}\s+\d{1,2}:\d{2}",
        r"\d{1,2}[/-]\d{1,2}[/-]\d{2,4}",
        r"\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{2,4}",
        r"(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s+\d{2,4}",
        r"\d{1,2}:\d{2}\s*(?:AM|PM|am|pm)",
        r"\d{1,2}:\d{2}\s*(?:IST|EST|PST|GMT|UTC)",
    ]

    # ------------------------------------------------------------------
    # Entry point
    # ------------------------------------------------------------------

    def parse(self, file_data: bytes, filename: str) -> DocumentContent:
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
        dispatch = {
            "txt":  self._parse_txt,
            "csv":  self._parse_csv,
            "xlsx": self._parse_excel,
            "xls":  self._parse_excel,
            "pdf":  self._parse_pdf,
            "docx": self._parse_docx,
            "pptx": self._parse_pptx,
            "msg":  self._parse_msg,
        }
        if ext not in dispatch:
            raise ValueError(f"Unsupported format: '{ext}'. Supported: {', '.join(dispatch)}")
        return dispatch[ext](file_data, filename)

    # ------------------------------------------------------------------
    # Text / CSV / Excel
    # ------------------------------------------------------------------

    def _parse_txt(self, data: bytes, filename: str) -> DocumentContent:
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("latin-1")

        paras = text.split("\n\n")
        chunk = max(1, len(paras) // 10)
        pages = {}
        for i in range(0, len(paras), chunk):
            pages[i // chunk + 1] = "\n\n".join(paras[i : i + chunk])

        return DocumentContent(
            text=text, total_pages=len(pages), page_contents=pages,
            metadata={"filename": filename}, file_type="txt",
        )

    def _parse_csv(self, data: bytes, filename: str) -> DocumentContent:
        import pandas as pd
        try:
            df = pd.read_csv(io.BytesIO(data))
        except Exception:
            df = pd.read_csv(io.BytesIO(data), encoding="latin-1")

        text = df.to_string(index=False)
        rows_per_page = 50
        pages = {}
        for i in range(0, len(df), rows_per_page):
            pages[i // rows_per_page + 1] = df.iloc[i : i + rows_per_page].to_string(index=False)

        return DocumentContent(
            text=text, total_pages=max(1, len(pages)), page_contents=pages,
            metadata={"filename": filename, "columns": list(df.columns), "rows": len(df)},
            file_type="csv",
        )

    def _parse_excel(self, data: bytes, filename: str) -> DocumentContent:
        import pandas as pd
        xf = pd.ExcelFile(io.BytesIO(data))
        all_text, pages = [], {}
        page_num = 1

        for sheet in xf.sheet_names:
            df = pd.read_excel(xf, sheet_name=sheet)
            block = f"=== Sheet: {sheet} ===\n{df.to_string(index=False)}"
            all_text.append(block)
            pages[page_num] = block
            page_num += 1

        return DocumentContent(
            text="\n\n".join(all_text), total_pages=len(pages), page_contents=pages,
            metadata={"filename": filename, "sheets": xf.sheet_names},
            file_type="excel",
        )

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def _parse_pdf(self, data: bytes, filename: str) -> DocumentContent:
        import pdfplumber
        pages, all_text = {}, []

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                pt = page.extract_text() or ""
                pages[i] = pt
                all_text.append(pt)

        images = self._pdf_page_images(data, total)
        return DocumentContent(
            text="\n\n".join(all_text), total_pages=total, page_contents=pages,
            metadata={"filename": filename, "page_count": total,
                      "image_count": len(images), "has_images": bool(images)},
            file_type="pdf", images=images,
        )

    def _pdf_page_images(self, data: bytes, total: int) -> List[ExtractedImage]:
        imgs = []
        try:
            from pdf2image import convert_from_bytes
            pages = convert_from_bytes(
                data, first_page=1, last_page=min(total, self.MAX_IMAGES), dpi=150
            )
            for n, img in enumerate(pages, 1):
                md = self.MAX_IMAGE_DIMENSION
                if img.width > md or img.height > md:
                    r = min(md / img.width, md / img.height)
                    img = img.resize((int(img.width * r), int(img.height * r)))
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                imgs.append(ExtractedImage(
                    image_base64=base64.b64encode(buf.getvalue()).decode(),
                    image_format="png", page_number=n, image_index=1,
                    description=f"Page {n}", width=img.width, height=img.height,
                ))
        except Exception:
            pass
        return imgs

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _build_text_pages(self, text: str) -> Dict[int, str]:
        estimated_pages = max(1, len(text.split()) // 500)
        chars_per_page = len(text) // estimated_pages if estimated_pages else len(text)
        return {i + 1: text[i * chars_per_page : (i + 1) * chars_per_page] for i in range(estimated_pages)}

    def _extract_text_from_document_xml(self, data: bytes) -> str:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            xml_bytes = archive.read("word/document.xml")
        xml_text = xml_bytes.decode("utf-8", errors="ignore")
        # Convert paragraph boundaries before stripping all tags.
        xml_text = re.sub(r"</w:p>", "\n", xml_text)
        xml_text = re.sub(r"<[^>]+>", "", xml_text)
        return unescape(xml_text).strip()

    def _best_effort_text_from_non_docx(self, data: bytes) -> str:
        for encoding in ("utf-8", "latin-1"):
            try:
                decoded = data.decode(encoding)
                cleaned = re.sub(r"\s+", " ", decoded).strip()
                if cleaned:
                    return cleaned
            except Exception:
                continue
        return ""

    def _parse_docx(self, data: bytes, filename: str) -> DocumentContent:
        from docx import Document

        try:
            doc = Document(io.BytesIO(data))
        except Exception as exc:
            # Some uploads are mis-labeled as .docx even though they are plain text.
            if not zipfile.is_zipfile(io.BytesIO(data)):
                fallback_text = self._best_effort_text_from_non_docx(data)
                if fallback_text and len(fallback_text) > 40:
                    pages = self._build_text_pages(fallback_text)
                    return DocumentContent(
                        text=fallback_text,
                        total_pages=len(pages),
                        page_contents=pages,
                        metadata={
                            "filename": filename,
                            "parser_note": "Used text fallback because file content is not a valid DOCX archive.",
                        },
                        file_type="docx",
                        images=[],
                    )
                raise ValueError(
                    "Invalid DOCX file. The file content is not a DOCX archive. "
                    "Please open it in Word/Google Docs and re-save as .docx."
                ) from exc

            # If the archive is valid but python-docx parsing fails, fallback to raw XML text extraction.
            try:
                fallback_text = self._extract_text_from_document_xml(data)
            except Exception:
                raise ValueError(
                    "Corrupted DOCX archive (ZIP central directory issue). "
                    "Please re-download or re-save the file and upload again."
                ) from exc

            if not fallback_text:
                raise ValueError(
                    "DOCX is readable but no text content was found in word/document.xml."
                ) from exc

            pages = self._build_text_pages(fallback_text)
            return DocumentContent(
                text=fallback_text,
                total_pages=len(pages),
                page_contents=pages,
                metadata={
                    "filename": filename,
                    "parser_note": "Used XML fallback because python-docx could not parse this file.",
                },
                file_type="docx",
                images=[],
            )

        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = [
            "\n".join(" | ".join(c.text for c in r.cells) for r in t.rows)
            for t in doc.tables
        ]
        text = "\n\n".join(paras + tables)
        pages = self._build_text_pages(text)

        images, count = [], 0
        try:
            for rel in doc.part.rels.values():
                if count >= self.MAX_IMAGES:
                    break
                if "image" in rel.reltype:
                    try:
                        blob = self._maybe_resize(rel.target_part.blob)
                        if len(blob) > 1000:
                            ct = rel.target_part.content_type
                            fmt = "jpeg" if "jpeg" in ct else "png"
                            images.append(ExtractedImage(
                                image_base64=base64.b64encode(blob).decode(),
                                image_format=fmt, page_number=1,
                                image_index=count + 1, description=f"Image {count + 1}",
                            ))
                            count += 1
                    except Exception:
                        pass
        except Exception:
            pass

        return DocumentContent(
            text=text, total_pages=len(pages), page_contents=pages,
            metadata={"filename": filename, "image_count": len(images),
                      "has_images": bool(images)},
            file_type="docx", images=images,
        )

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _parse_pptx(self, data: bytes, filename: str) -> DocumentContent:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(data))
        pages, all_text, images, count = {}, [], [], 0

        for snum, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_text.append(" | ".join(c.text for c in row.cells))
                if count < self.MAX_IMAGES:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            blob = self._maybe_resize(shape.image.blob)
                            if len(blob) > 1000:
                                ext = getattr(shape.image, "ext", "png").lower()
                                images.append(ExtractedImage(
                                    image_base64=base64.b64encode(blob).decode(),
                                    image_format="jpeg" if ext == "jpg" else ext,
                                    page_number=snum, image_index=count + 1,
                                    description=f"Slide {snum} image",
                                ))
                                count += 1
                    except Exception:
                        pass

            content = "\n".join(slide_text)
            pages[snum] = f"Slide {snum}:\n{content}"
            all_text.append(content)

        return DocumentContent(
            text="\n\n".join(all_text), total_pages=len(prs.slides), page_contents=pages,
            metadata={"filename": filename, "slide_count": len(prs.slides),
                      "image_count": len(images), "has_images": bool(images)},
            file_type="pptx", images=images,
        )

    # ------------------------------------------------------------------
    # MSG (Outlook e-mail)
    # ------------------------------------------------------------------

    def _parse_msg(self, data: bytes, filename: str) -> DocumentContent:
        try:
            import extract_msg
        except ImportError:
            raise ImportError("Install extract-msg:  pip install extract-msg")

        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(data)
            tmp_path = tmp.name

        try:
            msg = extract_msg.Message(tmp_path)
            sender_name = str(getattr(msg, "senderName", "") or "").strip() or "Unknown"
            sender_email = str(getattr(msg, "senderEmail", "") or "").strip()
            if "@" in sender_name and not sender_email:
                sender_email, sender_name = sender_name, sender_name.split("@")[0].replace(".", " ").title()
            full_sender = f"{sender_name} <{sender_email}>" if sender_email and sender_email != sender_name else sender_name

            to_list = [r.strip() for r in str(msg.to or "").replace(",", ";").split(";") if r.strip()]
            cc_list = [r.strip() for r in str(msg.cc or "").replace(",", ";").split(";") if r.strip()]

            date_str = "Not specified"
            if msg.date:
                try:
                    date_str = msg.date.strftime("%Y-%m-%d %H:%M:%S")
                except Exception:
                    date_str = str(msg.date)

            subject = str(msg.subject or "No Subject").strip()
            body = str(msg.body or "")
            if not body.strip() and getattr(msg, "htmlBody", None):
                body = self._strip_html(str(msg.htmlBody))

            original_body = body
            body = self._remove_disclaimer(body)

            main_links = self._extract_links(body)
            main_timelines = self._extract_timelines(body)
            main_action = self._detect_action(body)

            attachments: List[EmailAttachment] = []
            for att in msg.attachments:
                a = self._process_attachment(att)
                if a:
                    attachments.append(a)

            chain_msgs = self._parse_chain(original_body)
            msg.close()
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

        main_email = EmailMessage(
            sender=full_sender, recipients=to_list, cc=cc_list,
            subject=subject, date=date_str, body=body,
            attachments=attachments, links=main_links,
            timelines=main_timelines, action_requested=main_action,
        )

        all_msgs = chain_msgs + [main_email]
        for i, m in enumerate(all_msgs, 1):
            m.email_number = i

        all_links = [l for m in all_msgs for l in m.links]
        all_timelines = [t for m in all_msgs for t in m.timelines]

        participants = list({
            full_sender, *to_list, *cc_list,
            *[m.sender for m in chain_msgs if m.sender],
            *[r for m in chain_msgs for r in m.recipients],
        })
        participants = [p for p in participants if p]

        flow = self._build_flow(all_msgs)

        chain = EmailChain(
            messages=all_msgs, total_attachments=len(attachments),
            attachment_details=attachments, subject=subject,
            participants=participants, all_links=all_links,
            all_timelines=all_timelines, conversation_flow=flow,
        )

        content_lines = self._build_llm_content(all_msgs, subject, participants, all_links, all_timelines, flow)
        full_text = "\n".join(content_lines)

        page_contents = {
            i: f"Email #{m.email_number}\nFrom: {m.sender}\nDate: {m.date}\n\n{m.body}"
            for i, m in enumerate(all_msgs, 1)
        }

        return DocumentContent(
            text=full_text, total_pages=len(all_msgs), page_contents=page_contents,
            metadata={
                "filename": filename, "subject": subject, "sender": full_sender,
                "recipients": to_list, "cc": cc_list, "date": date_str,
                "participants": participants, "total_attachments": len(attachments),
                "attachment_names": [a.filename for a in attachments],
                "email_count": len(all_msgs), "is_email": True, "email_chain": chain,
                "all_links": [(l.url, l.link_type) for l in all_links],
                "all_timelines": [(t.original_text, t.event_type) for t in all_timelines],
                "conversation_flow": flow,
            },
            file_type="msg",
        )

    # ------------------------------------------------------------------
    # MSG helpers
    # ------------------------------------------------------------------

    def _extract_links(self, text: str) -> List[ExtractedLink]:
        results = []
        for m in self._URL_RE.finditer(text):
            url = m.group()
            ctx = text[max(0, m.start() - 50) : min(len(text), m.end() + 50)].replace("\n", " ")
            ul = url.lower()
            if any(k in ul for k in ("uat", "test", "staging")):
                ltype = "uat_link"
            elif any(k in ul for k in ("prod", "live")):
                ltype = "production_link"
            elif any(k in ul for k in ("sharepoint", "drive")):
                ltype = "document_link"
            elif any(k in ul for k in ("meet", "teams", "zoom")):
                ltype = "meeting_link"
            else:
                ltype = "url"
            results.append(ExtractedLink(url=url, context=ctx.strip(), link_type=ltype))
        return results

    def _extract_timelines(self, text: str) -> List[ExtractedTimeline]:
        results, seen = [], set()
        for pat in self._DATE_PATTERNS:
            for m in re.finditer(pat, text, re.IGNORECASE):
                raw = m.group().strip()
                if raw.lower() in seen:
                    continue
                seen.add(raw.lower())
                ctx = text[max(0, m.start() - 60) : min(len(text), m.end() + 60)].replace("\n", " ")
                cl = ctx.lower()
                if any(w in cl for w in ("meeting", "call", "walkthrough", "demo")):
                    etype = "meeting"
                elif any(w in cl for w in ("deadline", "due", "by", "before")):
                    etype = "deadline"
                elif any(w in cl for w in ("deploy", "release", "push", "uat")):
                    etype = "deployment"
                else:
                    etype = "general"
                results.append(ExtractedTimeline(original_text=raw, context=ctx.strip(), event_type=etype))
        return results

    def _detect_action(self, text: str) -> str:
        pats = [
            (r"(?:can you|could you|please|kindly)\s+([^.?!]{10,80})", "request"),
            (r"(?:sharing|attaching|attached|providing|here is)\s+([^.!]{10,80})", "sharing"),
            (r"(?:feedback|review|comments)\s+(?:on|for)\s+([^.!]{10,60})", "feedback_request"),
        ]
        actions = []
        for pat, atype in pats:
            for m in re.findall(pat, text, re.IGNORECASE):
                s = (m if isinstance(m, str) else " ".join(m)).strip()[:80]
                actions.append(f"[{atype}] {s}")
                if len(actions) >= 2:
                    break
        return "; ".join(actions[:2])

    def _build_flow(self, msgs: List[EmailMessage]) -> str:
        if len(msgs) <= 1:
            return ""
        parts = []
        for i, m in enumerate(msgs, 1):
            short = (m.sender.split("<")[0].strip() if "<" in m.sender else m.sender).split("@")[0]
            short = (short[:12] + "…") if len(short) > 12 else short
            parts.append(f"{i}. {short}")
        return " → ".join(parts)

    def _build_llm_content(self, msgs, subject, participants, links, timelines, flow) -> List[str]:
        lines = [
            "=" * 70,
            "EMAIL THREAD (CHRONOLOGICAL ORDER - Oldest to Newest)",
            "=" * 70,
            f"Subject: {subject}",
            f"Total Emails: {len(msgs)}",
            f"Participants: {', '.join(list(participants)[:5])}",
            "",
        ]
        if flow:
            lines += [f"CONVERSATION FLOW: {flow}", ""]
        if links:
            lines += ["-" * 50, f"ALL LINKS/URLs ({len(links)}):"]
            for l in links:
                lines.append(f"  • [{l.link_type}] {l.url}")
            lines.append("")
        if timelines:
            lines += ["-" * 50, f"ALL DATES/TIMES ({len(timelines)}):"]
            for t in timelines:
                lines.append(f"  • [{t.event_type}] {t.original_text}")
            lines.append("")
        lines += ["=" * 70, "DETAILED EMAIL SEQUENCE", "=" * 70]
        for m in msgs:
            lines += [
                "",
                f"--- EMAIL #{m.email_number} of {len(msgs)} ---",
                f"From: {m.sender}",
                f"To: {', '.join(m.recipients[:3])}",
                f"Date: {m.date}",
                f"Subject: {m.subject}",
                "",
                "MESSAGE:",
                m.body or "[No message body]",
            ]
            if m.links:
                lines.append(f"\nLINKS: {', '.join(l.url for l in m.links)}")
            if m.timelines:
                lines.append(f"\nDATES MENTIONED: {', '.join(t.original_text for t in m.timelines)}")
            if m.attachments:
                lines.append(f"\nATTACHMENTS: {', '.join(a.filename for a in m.attachments)}")
            lines += ["", "-" * 70]
        return lines

    def _parse_chain(self, body: str) -> List[EmailMessage]:
        if not body:
            return []
        pat = (
            r"(?:^|\n)\s*From:\s*([^\n]+)\s*\n"
            r"\s*(?:Sent|Date):\s*([^\n]+)\s*\n"
            r"\s*To:\s*([^\n]+)\s*\n"
            r"(?:\s*Cc:\s*([^\n]+)\s*\n)?"
            r"\s*Subject:\s*([^\n]+)\s*\n"
        )
        matches = list(re.finditer(pat, body, re.IGNORECASE | re.MULTILINE))
        msgs = []
        for i, m in enumerate(matches):
            try:
                sender = (m.group(1) or "").strip()
                date = (m.group(2) or "").strip()
                to = (m.group(3) or "").strip()
                subject = (m.group(5) or "").strip()
                start = m.end()
                end = matches[i + 1].start() if i + 1 < len(matches) else len(body)
                eb = self._remove_disclaimer(body[start:end].strip())
                if eb and len(eb) > 20:
                    msgs.append(EmailMessage(
                        sender=sender, recipients=[r.strip() for r in to.split(";") if r.strip()],
                        cc=[], subject=subject, date=date, body=eb[:5000],
                        links=self._extract_links(eb), timelines=self._extract_timelines(eb),
                        action_requested=self._detect_action(eb),
                    ))
            except Exception:
                continue
        msgs.reverse()   # oldest first
        return msgs

    def _process_attachment(self, att) -> Optional[EmailAttachment]:
        try:
            filename = str(getattr(att, "longFilename", None) or getattr(att, "shortFilename", None) or "unnamed")
            raw = getattr(att, "data", None)
            size = len(raw) if raw else 0
            if size == 0:
                return None
            ct = str(getattr(att, "mimetype", None) or "application/octet-stream")
            return EmailAttachment(filename=filename, content_type=ct, size_bytes=size)
        except Exception:
            return None

    def _strip_html(self, html: str) -> str:
        text = re.sub(r"<head[^>]*>.*?</head>", "", html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
        text = re.sub(r"</p>", "\n", text, flags=re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        return re.sub(r"[ \t]+", " ", text.replace("&nbsp;", " ").replace("&amp;", "&")).strip()

    def _remove_disclaimer(self, text: str) -> str:
        for pat in [
            r"Disclaimer:\s*This email communication.*?(?:write on|please write on).*",
            r"DISCLAIMER:.*?(?:prohibited|notify|delete).*?$",
        ]:
            try:
                text = re.sub(pat, "", text, flags=re.IGNORECASE | re.DOTALL | re.MULTILINE)
            except Exception:
                pass
        return text.strip()

    def _maybe_resize(self, data: bytes) -> bytes:
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(data))
            md = self.MAX_IMAGE_DIMENSION
            if img.width <= md and img.height <= md:
                return data
            r = min(md / img.width, md / img.height)
            img = img.resize((int(img.width * r), int(img.height * r)), Image.Resampling.LANCZOS)
            buf = io.BytesIO()
            fmt = img.format or "PNG"
            if img.mode == "RGBA" and fmt == "JPEG":
                img = img.convert("RGB")
            img.save(buf, format=fmt if fmt in ("PNG", "JPEG", "GIF") else "PNG")
            return buf.getvalue()
        except Exception:
            return data

    # ------------------------------------------------------------------
    # Range selection
    # ------------------------------------------------------------------

    def get_content_by_range(self, doc: DocumentContent, page_range: Optional[str] = None) -> str:
        if not page_range or not page_range.strip():
            return doc.text
        pages = self._parse_range(page_range, doc.total_pages)
        return "\n\n".join(doc.page_contents[p] for p in sorted(pages) if p in doc.page_contents) or doc.text

    def get_images_by_range(self, doc: DocumentContent, page_range: Optional[str] = None) -> List[ExtractedImage]:
        if not doc.images:
            return []
        if not page_range or not page_range.strip():
            return doc.images
        pages = set(self._parse_range(page_range, doc.total_pages))
        return [img for img in doc.images if img.page_number in pages]

    def _parse_range(self, s: str, max_val: int) -> List[int]:
        result = []
        for part in s.replace(" ", "").split(","):
            if "-" in part:
                try:
                    a, b = map(int, part.split("-", 1))
                    result.extend(range(max(1, a), min(b + 1, max_val + 1)))
                except ValueError:
                    pass
            else:
                try:
                    n = int(part)
                    if 1 <= n <= max_val:
                        result.append(n)
                except ValueError:
                    pass
        return sorted(set(result))
